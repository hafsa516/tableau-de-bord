from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from utils.chart_generator import create_chart


def get_logo_image():
    """Récupère l'image du logo"""
    logo_paths = [
        "logo_segula.png",
        "logo_segula.jpg",
        "logo_segula.jpeg",
        "logo.png",
    ]
    
    for path in logo_paths:
        if os.path.exists(path):
            try:
                with open(path, "rb") as f:
                    return BytesIO(f.read())
            except:
                continue
    return None


def create_pptx(pivot_table, output_path):
    """
    Crée un PowerPoint avec les graphiques
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ============================================================
    # SLIDE 1 : Logo Segula
    # ============================================================
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Logo
    logo_buffer = get_logo_image()
    if logo_buffer:
        try:
            slide.shapes.add_picture(logo_buffer, Inches(0), Inches(0),
                                     width=prs.slide_width, height=prs.slide_height)
        except:
            pass
    
    # ============================================================
    # SLIDE 2 : Graphique global
    # ============================================================
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "Répartition des statuts - Tous les projets"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = 1
    
    # Préparer les données
    if 'Total' in pivot_table.index:
        pivot_clean = pivot_table.drop('Total')
    else:
        pivot_clean = pivot_table
    
    if 'Total' in pivot_clean.columns:
        pivot_clean = pivot_clean.drop('Total', axis=1)
    
    categories = [str(col) for col in pivot_clean.columns]
    total_values = []
    for col in categories:
        val = pivot_clean[col].sum()
        total_values.append(int(val) if pd.notna(val) else 0)
    
    # Créer le graphique
    chart_buffer = create_chart(total_values, "Répartition des statuts - Tous les projets", categories)
    
    # Ajouter le graphique à la diapositive
    slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.2), width=Inches(12))
    
    # ============================================================
    # SLIDES SUIVANTS : Un graphique par projet
    # ============================================================
    if 'Total' in pivot_table.index:
        pivot_projects = pivot_table.drop('Total')
    else:
        pivot_projects = pivot_table
    
    if 'Total' in pivot_projects.columns:
        pivot_projects = pivot_projects.drop('Total', axis=1)
    
    for project_name in pivot_projects.index:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title.text_frame
        title_frame.text = f"Répartition des statuts - {project_name}"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = 1
        
        # Préparer les données
        project_data = pivot_projects.loc[project_name]
        categories = [str(col) for col in project_data.index]
        values = []
        for col in categories:
            values.append(int(project_data[col]) if pd.notna(project_data[col]) else 0)
        
        # Créer le graphique
        chart_buffer = create_chart(values, f"Répartition des statuts - {project_name}", categories)
        
        # Ajouter le graphique à la diapositive
        slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.2), width=Inches(12))
    
    # Sauvegarder
    prs.save(output_path)
    return output_path
